from flask import Flask, request, jsonify, send_file, render_template
from flask_cors import CORS
import os
import json
import uuid
import logging
from datetime import datetime
from werkzeug.utils import secure_filename

# NLP Libraries
from langdetect import detect, DetectorFactory
from deep_translator import GoogleTranslator
from gtts import gTTS

# Document Processing Libraries
from docx import Document
from pdfminer.high_level import extract_text
import PyPDF2
from openpyxl import load_workbook
from pptx import Presentation

# Set seed for consistent language detection
DetectorFactory.seed = 0

app = Flask(__name__)
CORS(app)

# Configuration
UPLOAD_FOLDER = 'uploads'
AUDIO_FOLDER = 'audio'
ALLOWED_EXTENSIONS = {'txt', 'pdf', 'docx', 'xlsx', 'pptx', 'doc'}
MAX_FILE_SIZE = 16 * 1024 * 1024  # 16MB

# Create directories
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(AUDIO_FOLDER, exist_ok=True)

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Language mappings for better NLP support
LANGUAGE_CODES = {
    'en': 'English',
    'es': 'Spanish', 
    'fr': 'French',
    'de': 'German',
    'it': 'Italian',
    'pt': 'Portuguese',
    'ru': 'Russian',
    'ja': 'Japanese',
    'ko': 'Korean',
    'zh': 'Chinese',
    'ar': 'Arabic',
    'hi': 'Hindi',
    'ta': 'Tamil',
    'te': 'Telugu',
    'bn': 'Bengali',
    'ur': 'Urdu',
    'th': 'Thai',
    'vi': 'Vietnamese',
    'nl': 'Dutch',
    'sv': 'Swedish',
    'da': 'Danish',
    'no': 'Norwegian',
    'fi': 'Finnish',
    'pl': 'Polish',
    'cs': 'Czech',
    'hu': 'Hungarian',
    'ro': 'Romanian',
    'bg': 'Bulgarian',
    'hr': 'Croatian',
    'sk': 'Slovak',
    'sl': 'Slovenian',
    'et': 'Estonian',
    'lv': 'Latvian',
    'lt': 'Lithuanian',
    'mt': 'Maltese',
    'ga': 'Irish',
    'cy': 'Welsh'
}

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_file(file_path):
    """Enhanced text extraction with better error handling"""
    try:
        file_extension = file_path.lower().split('.')[-1]
        
        if file_extension == 'pdf':
            # Try pdfminer first, fallback to PyPDF2
            try:
                text = extract_text(file_path)
                if not text.strip():
                    with open(file_path, 'rb') as file:
                        pdf_reader = PyPDF2.PdfReader(file)
                        text = ""
                        for page in pdf_reader.pages:
                            text += page.extract_text()
                return text
            except Exception as e:
                logger.error(f"PDF extraction error: {e}")
                return ""
                
        elif file_extension == 'docx':
            doc = Document(file_path)
            text = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += '\n' + cell.text
            return text
            
        elif file_extension == 'xlsx':
            workbook = load_workbook(file_path)
            text = ""
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                for row in sheet.iter_rows(values_only=True):
                    row_text = ' '.join([str(cell) for cell in row if cell is not None])
                    if row_text.strip():
                        text += row_text + '\n'
            return text
            
        elif file_extension == 'pptx':
            presentation = Presentation(file_path)
            text = ""
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + '\n'
            return text
            
        elif file_extension == 'txt':
            encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        return f.read()
                except UnicodeDecodeError:
                    continue
            return ""
            
        else:
            return ""
            
    except Exception as e:
        logger.error(f"Text extraction error for {file_path}: {e}")
        return ""

def detect_language_with_confidence(text):
    """Enhanced language detection with confidence scoring"""
    try:
        if not text or len(text.strip()) < 10:
            return 'en', 0.5
            
        # Clean text for better detection
        clean_text = ' '.join(text.split()[:1000])  # Use first 1000 words
        
        detected_lang = detect(clean_text)
        
        # Simple confidence estimation based on text length and character patterns
        confidence = min(0.9, 0.5 + (len(clean_text) / 1000) * 0.4)
        
        return detected_lang, confidence
        
    except Exception as e:
        logger.error(f"Language detection error: {e}")
        return 'en', 0.3

def translate_text_chunked(text, src_lang, dest_lang, max_chunk_size=4500):
    """Translate text in chunks to handle large documents using deep-translator"""
    try:
        # Handle auto-detection
        if src_lang == 'auto':
            src_lang = detect(text)
        
        # Create translator instance
        translator = GoogleTranslator(source=src_lang, target=dest_lang)
        
        if len(text) <= max_chunk_size:
            return translator.translate(text)
        
        # Split text into chunks
        chunks = []
        sentences = text.split('. ')
        current_chunk = ""
        
        for sentence in sentences:
            if len(current_chunk + sentence) < max_chunk_size:
                current_chunk += sentence + '. '
            else:
                if current_chunk:
                    chunks.append(current_chunk.strip())
                current_chunk = sentence + '. '
        
        if current_chunk:
            chunks.append(current_chunk.strip())
        
        # Translate each chunk
        translated_chunks = []
        for chunk in chunks:
            try:
                translated_chunk = translator.translate(chunk)
                translated_chunks.append(translated_chunk)
            except Exception as e:
                logger.error(f"Translation error for chunk: {e}")
                translated_chunks.append(chunk)  # Keep original if translation fails
        
        return ' '.join(translated_chunks)
        
    except Exception as e:
        logger.error(f"Translation error: {e}")
        return f"Translation failed: {str(e)}. Original text: {text[:200]}..."

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/api/languages')
def get_languages():
    """API endpoint to get supported languages"""
    return jsonify(LANGUAGE_CODES)

@app.route('/api/upload', methods=['POST'])
def upload_file():
    """Handle file upload and text extraction"""
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file uploaded'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        if not allowed_file(file.filename):
            return jsonify({'error': 'File type not supported'}), 400
        
        # Generate unique filename
        file_id = str(uuid.uuid4())
        filename = secure_filename(file.filename)
        file_path = os.path.join(UPLOAD_FOLDER, f"{file_id}_{filename}")
        
        file.save(file_path)
        
        # Extract text
        extracted_text = extract_text_from_file(file_path)
        
        if not extracted_text.strip():
            return jsonify({'error': 'No text could be extracted from the file'}), 400
        
        # Detect language
        detected_lang, confidence = detect_language_with_confidence(extracted_text)
        
        # Clean up uploaded file
        os.remove(file_path)
        
        return jsonify({
            'success': True,
            'text': extracted_text[:5000],  # Limit preview text
            'full_text': extracted_text,
            'detected_language': detected_lang,
            'language_name': LANGUAGE_CODES.get(detected_lang, 'Unknown'),
            'confidence': confidence,
            'word_count': len(extracted_text.split()),
            'char_count': len(extracted_text)
        })
        
    except Exception as e:
        logger.error(f"Upload error: {e}")
        return jsonify({'error': f'Upload failed: {str(e)}'}), 500

@app.route('/api/translate', methods=['POST'])
def translate_text():
    """Translate text to target language"""
    try:
        data = request.get_json()
        text = data.get('text', '')
        src_lang = data.get('source_language', 'auto')
        dest_lang = data.get('target_language', 'en')
        
        if not text:
            return jsonify({'error': 'No text provided'}), 400
        
        # Translate text
        translated_text = translate_text_chunked(text, src_lang, dest_lang)
        
        return jsonify({
            'success': True,
            'translated_text': translated_text,
            'source_language': src_lang,
            'target_language': dest_lang,
            'word_count': len(translated_text.split())
        })
        
    except Exception as e:
        logger.error(f"Translation error: {e}")
        return jsonify({'error': f'Translation failed: {str(e)}'}), 500

@app.route('/api/tts', methods=['POST'])
def text_to_speech():
    """Convert text to speech"""
    try:
        data = request.get_json()
        text = data.get('text', '')
        lang = data.get('language', 'en')
        
        if not text:
            return jsonify({'error': 'No text provided'}), 400
        
        # Limit text length for TTS
        if len(text) > 5000:
            text = text[:5000] + "..."
        
        # Generate unique filename
        audio_id = str(uuid.uuid4())
        audio_filename = f"tts_{audio_id}.mp3"
        audio_path = os.path.join(AUDIO_FOLDER, audio_filename)
        
        # Generate TTS
        tts = gTTS(text=text, lang=lang, slow=False)
        tts.save(audio_path)
        
        return jsonify({
            'success': True,
            'audio_url': f'/api/audio/{audio_filename}',
            'audio_id': audio_id
        })
        
    except Exception as e:
        logger.error(f"TTS error: {e}")
        return jsonify({'error': f'TTS generation failed: {str(e)}'}), 500

@app.route('/api/audio/<filename>')
def serve_audio(filename):
    """Serve audio files"""
    try:
        audio_path = os.path.join(AUDIO_FOLDER, filename)
        if os.path.exists(audio_path):
            return send_file(audio_path, mimetype='audio/mpeg')
        else:
            return jsonify({'error': 'Audio file not found'}), 404
    except Exception as e:
        logger.error(f"Audio serving error: {e}")
        return jsonify({'error': 'Failed to serve audio'}), 500

@app.route('/api/download/text', methods=['POST'])
def download_text():
    """Download translated text as file"""
    try:
        data = request.get_json()
        text = data.get('text', '')
        filename = data.get('filename', 'translated_text.txt')
        
        if not text:
            return jsonify({'error': 'No text provided'}), 400
        
        # Create temporary file
        temp_id = str(uuid.uuid4())
        temp_path = os.path.join(UPLOAD_FOLDER, f"temp_{temp_id}.txt")
        
        with open(temp_path, 'w', encoding='utf-8') as f:
            f.write(text)
        
        return send_file(temp_path, as_attachment=True, download_name=filename)
        
    except Exception as e:
        logger.error(f"Download error: {e}")
        return jsonify({'error': 'Download failed'}), 500

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)
